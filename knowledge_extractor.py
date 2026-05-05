import os
import re
import json
import jieba
import jieba.posseg as pseg
import networkx as nx
from collections import defaultdict, Counter
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


class DocumentParser:
    def __init__(self):
        self.stop_words = self._load_stop_words()
        jieba.initialize()

    def _load_stop_words(self):
        stop_words = set([
            '的', '了', '在', '是', '我', '有', '和', '就', '不', '人', '都', '一', '一个',
            '上', '也', '很', '到', '说', '要', '去', '你', '会', '着', '没有', '看', '好',
            '自己', '这', '那', '他', '她', '它', '们', '这个', '那个', '什么', '怎么', '为什么',
            'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with',
            'is', 'was', 'are', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'do', 'does',
            'did', 'will', 'would', 'could', 'should', 'may', 'might', 'must', 'can'
        ])
        return stop_words

    def parse_pdf(self, file_path):
        text = ""
        with open(file_path, 'rb') as f:
            reader = PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text

    def parse_docx(self, file_path):
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text

    def parse_pptx(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def parse_document(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            return self.parse_pdf(file_path)
        elif ext == '.docx':
            return self.parse_docx(file_path)
        elif ext == '.pptx':
            return self.parse_pptx(file_path)
        else:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()


class KnowledgeGraphGenerator:
    def __init__(self):
        self.parser = DocumentParser()
        self.documents = {}
        self.sentence_map = {}

    def extract_entities(self, text):
        words = pseg.cut(text)
        entities = []
        for word, flag in words:
            if word not in self.parser.stop_words and len(word) >= 2:
                if flag.startswith('n') or flag.startswith('nr') or flag.startswith('ns') or flag.startswith('nt'):
                    entities.append(word)
        return entities

    def extract_sentences(self, text):
        sentences = re.split(r'[。！？.!?\n]+', text)
        return [s.strip() for s in sentences if s.strip()]

    def build_graph(self, file_paths):
        G = nx.Graph()
        all_entities = []
        doc_id = 0

        for file_path in file_paths:
            doc_name = os.path.basename(file_path)
            text = self.parser.parse_document(file_path)
            self.documents[doc_id] = {'name': doc_name, 'text': text}

            sentences = self.extract_sentences(text)
            entities_in_doc = []

            for sent_idx, sentence in enumerate(sentences):
                entities = self.extract_entities(sentence)
                entities_in_doc.extend(entities)

                for entity in entities:
                    key = f"{doc_id}_{entity}"
                    if key not in self.sentence_map:
                        self.sentence_map[key] = []
                    self.sentence_map[key].append({
                        'sentence': sentence,
                        'doc_id': doc_id,
                        'doc_name': doc_name,
                        'sent_idx': sent_idx
                    })

                for i in range(len(entities)):
                    for j in range(i + 1, len(entities)):
                        if G.has_edge(entities[i], entities[j]):
                            G[entities[i]][entities[j]]['weight'] += 1
                        else:
                            G.add_edge(entities[i], entities[j], weight=1)

            all_entities.extend(entities_in_doc)
            doc_id += 1

        entity_counts = Counter(all_entities)
        for entity, count in entity_counts.items():
            G.add_node(entity, size=min(count * 5, 50), count=count)

        return G

    def graph_to_json(self, G):
        nodes = []
        edges = []

        for node, attrs in G.nodes(data=True):
            nodes.append({
                'id': node,
                'label': node,
                'size': attrs.get('size', 10),
                'count': attrs.get('count', 1)
            })

        for u, v, attrs in G.edges(data=True):
            edges.append({
                'source': u,
                'target': v,
                'weight': attrs.get('weight', 1)
            })

        return {'nodes': nodes, 'edges': edges}

    def get_entity_context(self, entity):
        contexts = []
        for key, sentences in self.sentence_map.items():
            if key.endswith(f'_{entity}'):
                contexts.extend(sentences)
        return contexts
